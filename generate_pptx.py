#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Bonian project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    PRIMARY_COLOR = RGBColor(102, 126, 234)  # #667eea
    SECONDARY_COLOR = RGBColor(118, 75, 162)  # #764ba2
    TEXT_COLOR = RGBColor(45, 55, 72)  # #2d3748
    LIGHT_TEXT = RGBColor(74, 85, 104)  # #4a5568
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add gradient background (simulated with shape)
    background = slide.shapes.add_shape(
        1,  # Rectangle
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY_COLOR
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Bonian"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "The Future of AI-Native Software Development"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.5))
    desc_frame = desc_box.text_frame
    desc_frame.text = "An IDE that turns UML diagrams into real, structured code"
    desc_para = desc_frame.paragraphs[0]
    desc_para.font.size = Pt(20)
    desc_para.font.color.rgb = RGBColor(255, 255, 255)
    desc_para.alignment = PP_ALIGN.CENTER
    
    # Team info
    team_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(2))
    team_frame = team_box.text_frame
    team_frame.word_wrap = True
    
    team_text = [
        "Team Members: Mohammed Mahjari, Ali Tawhari, Hatem Shwaiy,",
        "Saud Ghazwani, Haitham Torabi, Tawfiq Moharaq",
        "",
        "Supervisor: Dr. Mshabab Alrizah",
        "Jazan University – College of Computer Science"
    ]
    
    for i, line in enumerate(team_text):
        if i > 0:
            team_frame.add_paragraph()
        p = team_frame.paragraphs[i]
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Introduction", PRIMARY_COLOR)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    points = [
        "Software development is changing fast with the rise of AI-powered coding tools.",
        "While tools like GitHub Copilot and Cursor help developers write code faster, they rely heavily on natural language prompts, which can be ambiguous and error-prone.",
        "Bonian aims to fix this by introducing a clear, structured, and visual way to generate software using UML diagrams."
    ]
    
    add_bullet_points(content_frame, points, LIGHT_TEXT)
    
    # Slide 3: The Problem
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "The Problem", PRIMARY_COLOR)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    p = content_frame.paragraphs[0]
    p.text = "Most AI code assistants focus on speed, not architecture."
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_COLOR
    p.space_after = Pt(12)
    
    content_frame.add_paragraph()
    p = content_frame.paragraphs[1]
    p.text = 'They often forget earlier design decisions — a problem we call the Architectural Deficit.'
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_TEXT
    p.space_after = Pt(18)
    
    content_frame.add_paragraph()
    p = content_frame.paragraphs[2]
    p.text = "This leads to:"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_TEXT
    p.space_after = Pt(12)
    
    problems = ["Inconsistent code", "Security flaws", "Technical debt and poor maintainability"]
    for problem in problems:
        content_frame.add_paragraph()
        p = content_frame.paragraphs[-1]
        p.text = problem
        p.level = 1
        p.font.size = Pt(18)
        p.font.color.rgb = LIGHT_TEXT
    
    # Add highlight box
    highlight_box = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(1), Inches(5.5), Inches(8), Inches(0.8)
    )
    highlight_box.fill.solid()
    highlight_box.fill.fore_color.rgb = RGBColor(240, 242, 255)
    highlight_box.line.color.rgb = PRIMARY_COLOR
    highlight_box.line.width = Pt(3)
    
    text_frame = highlight_box.text_frame
    text_frame.text = "We believe the issue isn't the AI model itself — it's the input method."
    text_frame.paragraphs[0].font.size = Pt(18)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Slide 4: Our Vision
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Our Vision", PRIMARY_COLOR)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(4))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    points = [
        "Bonian redefines how developers interact with AI.",
        "Instead of typing prompts, developers design UML diagrams — and Bonian transforms them into executable code.",
        "This approach combines AI speed with engineering discipline."
    ]
    
    add_bullet_points(content_frame, points, LIGHT_TEXT)
    
    # Add highlight box
    highlight_box = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(1.5), Inches(5.5), Inches(7), Inches(0.8)
    )
    highlight_box.fill.solid()
    highlight_box.fill.fore_color.rgb = RGBColor(240, 242, 255)
    highlight_box.line.color.rgb = PRIMARY_COLOR
    highlight_box.line.width = Pt(3)
    
    text_frame = highlight_box.text_frame
    text_frame.text = "From prompt engineering to software architecture — that's Bonian."
    text_frame.paragraphs[0].font.size = Pt(18)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Slide 5: Objectives
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Objectives", PRIMARY_COLOR)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    points = [
        "Build an AI-native IDE that converts UML diagrams to code.",
        "Replace ambiguous text input with formal, visual modeling.",
        "Ensure architectural consistency in generated code.",
        "Improve developer productivity and code quality.",
        "Shift the developer's role from \"prompt user\" to system architect."
    ]
    
    add_bullet_points(content_frame, points, LIGHT_TEXT)
    
    # Slide 6: Project Scope
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Project Scope", PRIMARY_COLOR)
    
    # In Scope box
    in_scope_box = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.8), Inches(2), Inches(4), Inches(4)
    )
    in_scope_box.fill.solid()
    in_scope_box.fill.fore_color.rgb = RGBColor(240, 255, 244)
    in_scope_box.line.color.rgb = RGBColor(72, 187, 120)
    in_scope_box.line.width = Pt(2)
    
    text_frame = in_scope_box.text_frame
    text_frame.margin_top = Inches(0.2)
    text_frame.margin_left = Inches(0.2)
    text_frame.margin_right = Inches(0.2)
    
    p = text_frame.paragraphs[0]
    p.text = "✓ In Scope"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(34, 84, 61)
    p.space_after = Pt(12)
    
    in_scope_items = [
        "Desktop IDE (Windows, macOS, Linux)",
        "UML input → TypeScript output",
        "Real-time AI interpretation of diagrams"
    ]
    
    for item in in_scope_items:
        text_frame.add_paragraph()
        p = text_frame.paragraphs[-1]
        p.text = item
        p.level = 1
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_TEXT
    
    # Out of Scope box
    out_scope_box = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(5.2), Inches(2), Inches(4), Inches(4)
    )
    out_scope_box.fill.solid()
    out_scope_box.fill.fore_color.rgb = RGBColor(255, 245, 245)
    out_scope_box.line.color.rgb = RGBColor(229, 62, 62)
    out_scope_box.line.width = Pt(2)
    
    text_frame = out_scope_box.text_frame
    text_frame.margin_top = Inches(0.2)
    text_frame.margin_left = Inches(0.2)
    text_frame.margin_right = Inches(0.2)
    
    p = text_frame.paragraphs[0]
    p.text = "✗ Out of Scope"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(116, 42, 42)
    p.space_after = Pt(12)
    
    out_scope_items = [
        "Hand-drawn sketches",
        "Reverse engineering (code → UML)",
        "Multiple AI model integrations (Phase 2)"
    ]
    
    for item in out_scope_items:
        text_frame.add_paragraph()
        p = text_frame.paragraphs[-1]
        p.text = item
        p.level = 1
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_TEXT
    
    # Slide 7: Key Limitations
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Key Limitations", PRIMARY_COLOR)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(3.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    points = [
        "Dependent on the accuracy of multimodal AI models.",
        "May misread complex or non-standard UML diagrams.",
        "Forking VS Code creates a maintenance burden but allows deep customization."
    ]
    
    add_bullet_points(content_frame, points, LIGHT_TEXT)
    
    # Add highlight box
    highlight_box = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(1), Inches(5.5), Inches(8), Inches(0.8)
    )
    highlight_box.fill.solid()
    highlight_box.fill.fore_color.rgb = RGBColor(240, 242, 255)
    highlight_box.line.color.rgb = PRIMARY_COLOR
    highlight_box.line.width = Pt(3)
    
    text_frame = highlight_box.text_frame
    text_frame.text = "Despite these challenges, Bonian opens the door to new innovation beyond standard extensions."
    text_frame.paragraphs[0].font.size = Pt(16)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Slide 8: Literature Review
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Literature Review", PRIMARY_COLOR)
    
    # Add table
    rows, cols = 5, 3
    left = Inches(1)
    top = Inches(2.2)
    width = Inches(8)
    height = Inches(3.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(3)
    table.columns[2].width = Inches(3)
    
    # Header row
    headers = ["Tool", "Approach", "Limitation"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(18)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Data rows
    data = [
        ["Cursor", "Conversational AI in IDE", "Text-based only"],
        ["Copilot", "Code suggestion using GPT", "No architecture awareness"],
        ["UML Generators", "Rule-based", "No AI reasoning"],
        ["Bonian", "Vision + UML + AI", "Architecture-first design"]
    ]
    
    for i, row_data in enumerate(data, start=1):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = LIGHT_TEXT
            
            if j == 0:
                paragraph.font.bold = True
            
            # Alternate row colors
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(247, 250, 252)
    
    # Add conclusion text
    conclusion_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.8))
    conclusion_frame = conclusion_box.text_frame
    p = conclusion_frame.paragraphs[0]
    p.text = "Bonian bridges the gap between formal modeling and AI reasoning."
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 9: Methodology
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Methodology", PRIMARY_COLOR)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    p = content_frame.paragraphs[0]
    p.text = "Process Model: Agile + Research-based"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.space_after = Pt(18)
    
    content_frame.add_paragraph()
    p = content_frame.paragraphs[1]
    p.text = "Technologies: VS Code OSS, UML, TypeScript, GPT-4V or Claude 3"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.space_after = Pt(18)
    
    content_frame.add_paragraph()
    p = content_frame.paragraphs[2]
    p.text = "Development Phases:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.space_after = Pt(12)
    
    phases = [
        "Research & Design",
        "AI Model Integration",
        "System Implementation",
        "Testing & Evaluation"
    ]
    
    for phase in phases:
        content_frame.add_paragraph()
        p = content_frame.paragraphs[-1]
        p.text = phase
        p.level = 1
        p.font.size = Pt(18)
        p.font.color.rgb = LIGHT_TEXT
    
    # Slide 10: System Design
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "System Design", PRIMARY_COLOR)
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "Bonian includes three main components:"
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_COLOR
    
    # Three component boxes
    components = [
        ("AI Engine", "Interprets UML diagrams"),
        ("Code Generator", "Produces executable code"),
        ("IDE Interface", "Enables interaction and visualization")
    ]
    
    box_width = Inches(2.5)
    box_height = Inches(2)
    start_x = Inches(1)
    start_y = Inches(3)
    spacing = Inches(0.3)
    
    for i, (title, desc) in enumerate(components):
        x = start_x + i * (box_width + spacing)
        
        box = slide.shapes.add_shape(
            1,  # Rectangle
            x, start_y, box_width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(240, 242, 255)
        box.line.fill.background()
        
        text_frame = box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(12)
        
        text_frame.add_paragraph()
        p = text_frame.paragraphs[1]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_TEXT
        p.alignment = PP_ALIGN.CENTER
    
    # Add highlight box
    highlight_box = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(1.5), Inches(5.5), Inches(7), Inches(0.8)
    )
    highlight_box.fill.solid()
    highlight_box.fill.fore_color.rgb = RGBColor(240, 242, 255)
    highlight_box.line.color.rgb = PRIMARY_COLOR
    highlight_box.line.width = Pt(3)
    
    text_frame = highlight_box.text_frame
    text_frame.text = "It combines the precision of UML with the intelligence of AI."
    text_frame.paragraphs[0].font.size = Pt(18)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Slide 11: Project Plan
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Project Plan", PRIMARY_COLOR)
    
    phases = [
        ("Phase 1", "Research & Literature Review"),
        ("Phase 2", "Requirements & Planning"),
        ("Phase 3", "System Design & Implementation"),
        ("Phase 4", "Testing & Final Report")
    ]
    
    box_width = Inches(4)
    box_height = Inches(1.2)
    
    for i, (phase_num, phase_desc) in enumerate(phases):
        row = i // 2
        col = i % 2
        
        x = Inches(1) + col * (box_width + Inches(0.4))
        y = Inches(2.5) + row * (box_height + Inches(0.3))
        
        box = slide.shapes.add_shape(
            1,  # Rectangle
            x, y, box_width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(247, 250, 252)
        box.line.color.rgb = PRIMARY_COLOR
        box.line.width = Pt(2)
        
        text_frame = box.text_frame
        text_frame.margin_top = Inches(0.15)
        text_frame.margin_left = Inches(0.2)
        
        p = text_frame.paragraphs[0]
        p.text = phase_num
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.space_after = Pt(6)
        
        text_frame.add_paragraph()
        p = text_frame.paragraphs[1]
        p.text = phase_desc
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_TEXT
    
    # Add note
    note_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
    note_frame = note_box.text_frame
    p = note_frame.paragraphs[0]
    p.text = "(You can show your Gantt chart or WBS diagram here.)"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = RGBColor(113, 128, 150)
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 12: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Conclusion", PRIMARY_COLOR)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(4))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    p = content_frame.paragraphs[0]
    p.text = "Bonian is more than a tool — it's a new paradigm for AI-assisted development."
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_COLOR
    p.space_after = Pt(18)
    
    content_frame.add_paragraph()
    p = content_frame.paragraphs[1]
    p.text = "It aims to:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.space_after = Pt(12)
    
    aims = [
        "Strengthen software architecture",
        "Reduce technical debt",
        "Empower developers with design-first thinking"
    ]
    
    for aim in aims:
        content_frame.add_paragraph()
        p = content_frame.paragraphs[-1]
        p.text = aim
        p.level = 1
        p.font.size = Pt(18)
        p.font.color.rgb = LIGHT_TEXT
    
    # Add final highlight box
    highlight_box = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(1.5), Inches(5.5), Inches(7), Inches(0.9)
    )
    highlight_box.fill.solid()
    highlight_box.fill.fore_color.rgb = RGBColor(240, 242, 255)
    highlight_box.line.color.rgb = PRIMARY_COLOR
    highlight_box.line.width = Pt(3)
    
    text_frame = highlight_box.text_frame
    text_frame.text = '"From prompt engineering to software architecture — that\'s Bonian."'
    text_frame.paragraphs[0].font.size = Pt(20)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Save presentation
    prs.save('Bonian_Presentation.pptx')
    print("✓ PowerPoint presentation created successfully: Bonian_Presentation.pptx")

def add_title(slide, title_text, color):
    """Add a title to a slide"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = color
    
    # Add underline shape
    line = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(1.4), Inches(1.2), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def add_bullet_points(text_frame, points, color):
    """Add bullet points to a text frame"""
    for i, point in enumerate(points):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = color
        p.space_after = Pt(12)

if __name__ == "__main__":
    create_presentation()
